import json
import re
import threading
import uuid
import networkx as nx
import dataiku
from flask import request, Response

# --- State ---
sessions = {}

# --- LLM ---
def get_project():
    client = dataiku.api_client()
    return client.get_default_project()

def list_available_llms():
    """List all LLM IDs available in this Dataiku project."""
    project = get_project()
    try:
        llms = project.list_llms()
        return [m['id'] for m in llms]
    except Exception:
        return []

def get_llm(model_id=None):
    project = get_project()
    if model_id:
        return project.get_llm(model_id)
    # Default: pick the first available model
    available = list_available_llms()
    return project.get_llm(available[0] if available else 'openai:MSOpenAI:gpt-4o')

llm = get_llm()

# --- Orion / MS Color Palette ---
ORION_COLORS = ['#216CA6', '#1A936F', '#7B2D8E', '#C5283D', '#2E86AB',
                '#5ba3d9', '#e07a2f', '#6C757D', '#0d7377', '#8B5CF6']

# --- Domain Templates ---
DOMAIN_TEMPLATES = {
    'communication': {
        'name': 'Communication Intelligence',
        'description': 'Analyze emails, messages, and team communications to map people, topics, decisions, and information flow.',
        'entity_types': ['Person', 'Team', 'Topic', 'Decision', 'Action Item', 'Project', 'Meeting'],
        'relationship_types': ['SENT_TO', 'CC_TO', 'DISCUSSED', 'DECIDED', 'ASSIGNED_TO', 'BELONGS_TO', 'ESCALATED_TO', 'FOLLOWED_UP', 'MENTIONED']
    },
    'process': {
        'name': 'Process Optimization',
        'description': 'Map processes, systems, teams, and data flows to identify bottlenecks, manual handoffs, and automation opportunities.',
        'entity_types': ['Process', 'Step', 'System', 'Team', 'Data', 'Tool', 'Bottleneck'],
        'relationship_types': ['FEEDS_INTO', 'TRIGGERS', 'OWNED_BY', 'USES', 'PRODUCES', 'DEPENDS_ON', 'BLOCKED_BY', 'HANDOFF_TO']
    },
    'code': {
        'name': 'Code Review',
        'description': 'Analyze code, documentation, or technical specs to map modules, functions, dependencies, and features.',
        'entity_types': ['Module', 'Function', 'Class', 'API', 'Feature', 'Dependency', 'Config'],
        'relationship_types': ['CALLS', 'IMPORTS', 'EXTENDS', 'IMPLEMENTS', 'DEPENDS_ON', 'EXPOSES', 'CONFIGURED_BY', 'TESTED_BY']
    },
    'workflow': {
        'name': 'Workflow Orchestration',
        'description': 'Map workflow stages, triggers, approvals, and handoffs to visualize orchestration and dependencies.',
        'entity_types': ['Stage', 'Task', 'Actor', 'Trigger', 'Condition', 'Output', 'Queue'],
        'relationship_types': ['TRANSITIONS_TO', 'TRIGGERED_BY', 'APPROVED_BY', 'PRODUCES', 'WAITS_FOR', 'ROUTES_TO', 'ESCALATES_TO', 'PARALLEL_WITH']
    },
    'legal': {
        'name': 'Legal Document',
        'description': 'Extract parties, obligations, rights, clauses, and conditions from legal documents.',
        'entity_types': ['Party', 'Obligation', 'Right', 'Clause', 'Term', 'Jurisdiction', 'Condition'],
        'relationship_types': ['OBLIGATED_TO', 'HAS_RIGHT', 'GOVERNED_BY', 'CONDITIONAL_ON', 'DEFINED_IN', 'SUPERSEDES', 'EFFECTIVE_FROM', 'SUBJECT_TO']
    },
    'contract': {
        'name': 'Contract & Agreement',
        'description': 'Map parties, terms, amounts, dates, deliverables, and SLAs from contracts and agreements.',
        'entity_types': ['Party', 'Term', 'Amount', 'Date', 'Deliverable', 'SLA', 'Penalty'],
        'relationship_types': ['AGREED_BY', 'PAYS_TO', 'DELIVERS_TO', 'DUE_BY', 'PENALIZED_IF', 'RENEWED_ON', 'TERMINABLE_BY', 'REFERENCES']
    },
    'data_analysis': {
        'name': 'Data Analysis',
        'description': 'Map datasets, metrics, dimensions, sources, transformations, and insights from analytical content.',
        'entity_types': ['Dataset', 'Metric', 'Dimension', 'Source', 'Insight', 'Trend', 'Segment'],
        'relationship_types': ['MEASURED_BY', 'SOURCED_FROM', 'CORRELATES_WITH', 'SEGMENTED_BY', 'DRIVES', 'COMPARED_TO', 'DERIVED_FROM', 'INDICATES']
    },
    'general': {
        'name': 'General',
        'description': None,  # Will use auto-inference
        'entity_types': None,
        'relationship_types': None
    }
}

# --- Prompts ---
ENTITY_VS_ATTRIBUTE_GUIDANCE = """
CRITICAL — Entity vs. Attribute distinction:
- An ENTITY is something that has its own identity and participates in relationships (e.g. Person, Organization, Project, System)
- An ATTRIBUTE is a property of an entity that describes it but does NOT need its own node (e.g. phone number, email address, employee ID, floor number, zip code, date of birth)
- DO NOT create entity types for attributes. Instead, attributes should be included in the entity's "description" field during extraction.
- Examples of what should be ENTITIES: Person, Team, Division, Title/Role, Building, Project, System
- Examples of what should be ATTRIBUTES (NOT entities): Phone Number, Email, Employee ID, Floor, Extension, Badge Number, Hire Date, Address
- Ask yourself: "Would this create a meaningful node in a graph that connects to other things?" If not, it's an attribute.
"""

SCHEMA_PROMPT_GENERAL = """You are a data analyst. Read the following text and define a focused ontology for a knowledge graph.

Your job:
1. Understand what this data is about (e.g. employee directory, financial report, research paper, etc.)
2. Define a SMALL set of entity types (3-8 types max) that capture the key concepts
3. Define a SMALL set of relationship types (3-10 types max) that capture how entities connect
4. Keep it focused — fewer, meaningful types are better than many granular ones
5. Identify fields that are ATTRIBUTES (not entities) so they can be excluded from the graph
{entity_attr_guidance}

Return a JSON object:
{{
  "description": "Brief description of what this data is about",
  "domain": "auto-detected domain name",
  "entity_types": ["Type1", "Type2", "Type3"],
  "relationship_types": ["REL_TYPE_1", "REL_TYPE_2"],
  "attribute_fields": ["field1", "field2"]
}}

Rules:
- Use clear, singular nouns for entity types (e.g. "Person" not "People")
- Use UPPER_SNAKE_CASE for relationship types (e.g. "HAS_TITLE" not "has title")
- Merge similar concepts (e.g. don't have both "Role" and "Title" — pick one)
- "attribute_fields" should list data fields that are properties/attributes and should NOT become graph nodes
- Return ONLY valid JSON, no markdown fences, no extra text.

Text (sample):
{text}
"""

SCHEMA_PROMPT_AUTO = """You are a data analyst. Read the following text and determine which analysis domain fits best.

Available domains:
{domain_list}
{entity_attr_guidance}

Your job:
1. Read the text sample and classify it into the BEST matching domain above
2. Use that domain's suggested entity types and relationship types as a starting point
3. Adapt the ontology to fit the actual data — add, remove, or rename types as needed (stay within 3-8 entity types, 3-10 relationship types)
4. Identify fields that are ATTRIBUTES (not entities) — these should be excluded from graph nodes
5. If no domain fits well, create a custom ontology from scratch

Return a JSON object:
{{
  "description": "Brief description of what this data is about",
  "domain": "detected domain key (e.g. 'communication', 'legal', 'general')",
  "entity_types": ["Type1", "Type2", "Type3"],
  "relationship_types": ["REL_TYPE_1", "REL_TYPE_2"],
  "attribute_fields": ["field1", "field2"]
}}

Rules:
- Use clear, singular nouns for entity types
- Use UPPER_SNAKE_CASE for relationship types
- "attribute_fields" should list data fields that are properties/attributes and should NOT become graph nodes (e.g. phone numbers, IDs, email addresses, floor numbers)
- Return ONLY valid JSON, no markdown fences, no extra text.

Text (sample):
{text}
"""

SCHEMA_PROMPT_DOMAIN = """You are a domain expert analyzing data for: {domain_name}

Domain context: {domain_desc}
Suggested entity types: {entity_types}
Suggested relationship types: {rel_types}
{entity_attr_guidance}

Read the following text and refine the ontology to fit this specific data.
You MUST use the suggested types as your starting point, but you may:
- Add 1-2 additional types if the data clearly needs them
- Remove types that don't appear in the data
- Keep the total within 3-8 entity types and 3-10 relationship types
- Identify fields that are ATTRIBUTES (not entities) — these should be excluded from graph nodes

Return a JSON object:
{{
  "description": "Brief description of what this data is about",
  "domain": "{domain_key}",
  "entity_types": ["Type1", "Type2", "Type3"],
  "relationship_types": ["REL_TYPE_1", "REL_TYPE_2"],
  "attribute_fields": ["field1", "field2"]
}}

Rules:
- Use clear, singular nouns for entity types
- Use UPPER_SNAKE_CASE for relationship types
- "attribute_fields" should list data fields that are properties/attributes and should NOT become graph nodes
- Return ONLY valid JSON, no markdown fences, no extra text.

Text (sample):
{text}
"""

EXTRACTION_PROMPT = """Extract entities and relationships from the text below using ONLY the provided ontology.

Ontology:
- Entity types: {entity_types}
- Relationship types: {relationship_types}
- Context: {schema_description}
- Attribute fields (DO NOT create nodes for these — include them in entity descriptions instead): {attribute_fields}

Return a JSON object:
{{
  "entities": [
    {{"name": "Entity Name", "type": "EntityType", "description": "brief description including relevant attributes"}}
  ],
  "relationships": [
    {{"source": "Entity Name", "target": "Entity Name", "relation": "RELATIONSHIP_TYPE", "description": "brief description"}}
  ]
}}

Rules:
- Use ONLY the entity types and relationship types listed above. Do not invent new ones.
- Extract meaningful entities that participate in relationships. Skip trivial or isolated data.
- Attribute fields (like phone numbers, IDs, email addresses, floor numbers) should be folded into the entity's "description" field, NOT created as separate nodes.
- For tabular/record data: the primary entity is the row subject (e.g. a Person). Other columns are either related entities (if they match an entity type) or attributes (include in description).
- Return ONLY valid JSON, no markdown fences, no extra text.

Text:
{text}
"""

DEDUP_PROMPT = """You are an entity resolution agent. Given a list of entity names extracted from text,
identify groups of names that refer to the SAME real-world entity (duplicates, abbreviations, typos, partial names).

Entity list:
{entities}

Return a JSON object mapping each variant name to its canonical (best/fullest) form:
{{
  "mappings": {{
    "variant_name": "Canonical Name",
    "another_variant": "Canonical Name"
  }}
}}

Rules:
- Only include names that need to be merged. Skip names that are already unique.
- Choose the most complete, properly spelled version as the canonical name.
- Examples: "J.P. Morgan" and "JPMorgan Chase" -> use "JPMorgan Chase". "Ted P." and "Ted Pick" -> use "Ted Pick".
- If unsure, do NOT merge. Only merge when clearly the same entity.
- Return ONLY valid JSON, no markdown fences, no extra text.
"""

def json_response(data, status=200):
    return Response(json.dumps(data), status=status, mimetype='application/json')

# --- Helper Functions ---
def clean_llm_json(raw):
    text = raw.strip()
    text = re.sub(r'^```(?:json)?\s*\n?', '', text)
    text = re.sub(r'\n?```\s*$', '', text)
    text = text.replace('\\n', '\n')
    match = re.search(r'\{[\s\S]*\}', text)
    if match:
        text = match.group(0)
    return text.strip()


def infer_schema(text, llm_inst=None, domain='auto'):
    """Agent 0: Analyze data and define a focused ontology."""
    _llm = llm_inst or llm
    sample = text[:2000]

    guidance = ENTITY_VS_ATTRIBUTE_GUIDANCE

    if domain != 'auto' and domain in DOMAIN_TEMPLATES and DOMAIN_TEMPLATES[domain]['entity_types']:
        # User selected a specific domain template
        tmpl = DOMAIN_TEMPLATES[domain]
        prompt = SCHEMA_PROMPT_DOMAIN.format(
            domain_name=tmpl['name'],
            domain_desc=tmpl['description'],
            entity_types=', '.join(tmpl['entity_types']),
            rel_types=', '.join(tmpl['relationship_types']),
            domain_key=domain,
            entity_attr_guidance=guidance,
            text=sample
        )
    elif domain == 'auto':
        # Auto-detect: give Agent 0 the full domain list to classify
        domain_list = ''
        for key, tmpl in DOMAIN_TEMPLATES.items():
            if key == 'general' or not tmpl['entity_types']:
                continue
            domain_list += '- ' + key + ' (' + tmpl['name'] + '): ' + tmpl['description'] + '\n'
            domain_list += '  Entity types: ' + ', '.join(tmpl['entity_types']) + '\n'
            domain_list += '  Relationship types: ' + ', '.join(tmpl['relationship_types']) + '\n'
        prompt = SCHEMA_PROMPT_AUTO.format(domain_list=domain_list, entity_attr_guidance=guidance, text=sample)
    else:
        # General / fallback
        prompt = SCHEMA_PROMPT_GENERAL.format(entity_attr_guidance=guidance, text=sample)

    completion = _llm.new_completion()
    completion.with_message(prompt)
    resp = completion.execute()
    cleaned = clean_llm_json(resp.text)
    return json.loads(cleaned)


def extract_entities(text, schema, llm_inst=None):
    _llm = llm_inst or llm
    completion = _llm.new_completion()
    attr_fields = schema.get('attribute_fields', [])
    prompt = EXTRACTION_PROMPT.format(
        text=text,
        entity_types=', '.join(schema['entity_types']),
        relationship_types=', '.join(schema['relationship_types']),
        schema_description=schema.get('description', ''),
        attribute_fields=', '.join(attr_fields) if attr_fields else 'none specified'
    )
    completion.with_message(prompt)
    resp = completion.execute()
    cleaned = clean_llm_json(resp.text)
    return json.loads(cleaned)


def deduplicate_entities(G, llm_inst=None):
    _llm = llm_inst or llm
    entity_names = list(G.nodes())
    if len(entity_names) < 2:
        return G

    completion = _llm.new_completion()
    completion.with_message(DEDUP_PROMPT.format(entities='\n'.join(entity_names)))
    resp = completion.execute()
    cleaned = clean_llm_json(resp.text)
    result = json.loads(cleaned)
    mappings = result.get('mappings', {})

    if not mappings:
        return G

    # Build new graph with merged entities
    new_G = nx.DiGraph()

    # Copy nodes, merging as needed
    for node, data in G.nodes(data=True):
        canonical = mappings.get(node, node)
        if canonical not in new_G:
            new_G.add_node(canonical, **data)
        else:
            # Keep the richer description
            existing = new_G.nodes[canonical]
            if len(data.get('description', '')) > len(existing.get('description', '')):
                existing['description'] = data['description']
            if existing.get('type', 'Unknown') == 'Unknown' and data.get('type', 'Unknown') != 'Unknown':
                existing['type'] = data['type']

    # Copy edges, remapping to canonical names
    for src, tgt, data in G.edges(data=True):
        new_src = mappings.get(src, src)
        new_tgt = mappings.get(tgt, tgt)
        if new_src == new_tgt:
            continue  # Skip self-loops created by merging
        if not new_G.has_edge(new_src, new_tgt):
            if new_src not in new_G:
                new_G.add_node(new_src, type='Unknown', description='')
            if new_tgt not in new_G:
                new_G.add_node(new_tgt, type='Unknown', description='')
            new_G.add_edge(new_src, new_tgt, **data)

    return new_G


def chunk_text(text, chunk_size=400, overlap=30):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            for sep in ['. ', '\n\n', '\n', ' ']:
                last_sep = text[start:end].rfind(sep)
                if last_sep > chunk_size * 0.5:
                    end = start + last_sep + len(sep)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap
    return chunks


def graph_to_json(G):
    type_color_map = {}
    nodes = []
    for node, data in G.nodes(data=True):
        node_type = data.get('type', 'Unknown')
        if node_type not in type_color_map:
            idx = len(type_color_map) % len(ORION_COLORS)
            type_color_map[node_type] = ORION_COLORS[idx]
        color = type_color_map[node_type]
        nodes.append({
            'id': node,
            'label': node,
            'fullName': node,
            'type': node_type,
            'description': data.get('description', ''),
            'color': color
        })
    edges = []
    for src, tgt, data in G.edges(data=True):
        edges.append({
            'from': src, 'to': tgt,
            'relation': data.get('relation', ''),
            'description': data.get('description', '')
        })
    return {'nodes': nodes, 'edges': edges, 'typeColors': type_color_map}


def get_full_context(G):
    lines = ['Knowledge Graph Summary:', '']
    by_type = {}
    for node, data in G.nodes(data=True):
        t = data.get('type', 'Unknown')
        by_type.setdefault(t, []).append((node, data))
    lines.append('ENTITIES:')
    for entity_type, nodes_list in by_type.items():
        lines.append('  ' + entity_type + ':')
        for name, data in nodes_list:
            lines.append('    - ' + name + ': ' + data.get('description', ''))
    lines.append('')
    lines.append('RELATIONSHIPS:')
    for src, tgt, data in G.edges(data=True):
        lines.append('  ' + src + ' --[' + data['relation'] + ']--> ' + tgt + ': ' + data.get('description', ''))
    return '\n'.join(lines)


def generate_summary(G):
    """Generate a one-line summary card from graph data."""
    node_count = G.number_of_nodes()
    edge_count = G.number_of_edges()
    types = {}
    for _, data in G.nodes(data=True):
        t = data.get('type', 'Unknown')
        types[t] = types.get(t, 0) + 1
    # Find top entity (most connections)
    top_entity = None
    max_degree = 0
    for node in G.nodes():
        deg = G.degree(node)
        if deg > max_degree:
            max_degree = deg
            top_entity = node
    type_str = ', '.join(str(v) + ' ' + k + ('s' if v > 1 else '') for k, v in
                         sorted(types.items(), key=lambda x: -x[1]))
    summary = 'Discovered ' + str(node_count) + ' entities across ' + str(len(types)) + ' types (' + type_str + ') with ' + str(edge_count) + ' connections.'
    if top_entity:
        summary += ' Most connected: ' + top_entity + ' (' + str(max_degree) + ' connections).'
    return summary


def generate_report(session, llm_inst=None):
    """Generate an analysis report from graph and source data."""
    _llm = llm_inst or llm
    ctx = get_full_context(session['graph'])
    source = session['source_text']
    if len(source) > 8000:
        source = source[:8000] + '\n[... truncated ...]'
    prompt = """You are an analyst generating a research report.
Based on the following entity and relationship data, write a detailed analysis report.

""" + ctx + """

Original source text:
""" + source + """

Write a report with these sections:
1. Executive Summary (2-3 sentences)
2. Key Entities & Their Roles
3. Key Relationships & Dynamics
4. Outlook & Implications

Be concise and professional."""

    completion = _llm.new_completion()
    completion.with_message(prompt)
    resp = completion.execute()
    return resp.text


# --- Background graph building ---
def build_graph_async(session_id, text, model_id=None, domain='auto'):
    session = sessions[session_id]
    session['status'] = 'analyzing'
    session['source_text'] = text
    llm_inst = get_llm(model_id) if model_id else llm
    session['llm_inst'] = llm_inst

    # Agent 0: Infer schema/ontology
    try:
        schema = infer_schema(text, llm_inst, domain)
        session['schema'] = schema
    except Exception as e:
        session.setdefault('errors', []).append('Schema: ' + str(e))
        schema = {'entity_types': ['Entity'], 'relationship_types': ['RELATED_TO'], 'description': ''}

    if session.get('stopped'):
        session['status'] = 'stopped'
        return

    session['status'] = 'building'
    chunks = chunk_text(text)
    total = len(chunks)
    session['total_chunks'] = total
    G = nx.DiGraph()

    for i, chunk in enumerate(chunks):
        if session.get('stopped'):
            session['status'] = 'stopped'
            session['graph'] = G
            session['graph_data'] = graph_to_json(G)
            return
        session['current_chunk'] = i + 1
        try:
            result = extract_entities(chunk, schema, llm_inst)
            for entity in result.get('entities', []):
                if entity['name'] not in G:
                    G.add_node(entity['name'], type=entity['type'],
                               description=entity['description'])
            for rel in result.get('relationships', []):
                if rel['source'] not in G:
                    G.add_node(rel['source'], type='Unknown', description='')
                if rel['target'] not in G:
                    G.add_node(rel['target'], type='Unknown', description='')
                if not G.has_edge(rel['source'], rel['target']):
                    G.add_edge(rel['source'], rel['target'],
                               relation=rel['relation'],
                               description=rel['description'])
            session['graph'] = G
            session['graph_data'] = graph_to_json(G)
        except Exception as e:
            session.setdefault('errors', []).append('Chunk ' + str(i+1) + ': ' + str(e))

    # Agent 2: Entity deduplication
    if session.get('stopped'):
        session['status'] = 'stopped'
        return
    if G.number_of_nodes() > 1:
        session['status'] = 'deduplicating'
        try:
            G = deduplicate_entities(G, llm_inst)
            session['graph'] = G
            session['graph_data'] = graph_to_json(G)
        except Exception as e:
            session.setdefault('errors', []).append('Dedup: ' + str(e))

    # Auto-generate summary
    try:
        summary = generate_summary(G)
        session['summary'] = summary
    except Exception as e:
        session.setdefault('errors', []).append('Summary: ' + str(e))

    # Auto-generate report
    if session.get('stopped'):
        session['status'] = 'stopped'
        return
    session['status'] = 'reporting'
    try:
        report = generate_report(session, llm_inst)
        session['report'] = report
    except Exception as e:
        session.setdefault('errors', []).append('Report: ' + str(e))

    session['status'] = 'done'


# --- File parsing (structure-preserving) ---
IMAGE_EXTENSIONS = ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp')

IMAGE_MIME_MAP = {
    'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
    'gif': 'image/gif', 'bmp': 'image/bmp', 'tiff': 'image/tiff',
    'tif': 'image/tiff', 'webp': 'image/webp',
}

OCR_SYSTEM_PROMPT = """You are Orion Document Vision — a world-class, enterprise-grade document intelligence engine. Your sole mission is PERFECT, LOSSLESS content extraction. You operate at the highest fidelity: every character, every pixel of meaning matters. Zero hallucination tolerance — if you cannot read it with certainty, mark it as [unclear]. Never invent, infer, or assume content that is not visually present."""

OCR_PROMPT = """Perform exhaustive content extraction on this document image. Your output will be fed into a knowledge graph pipeline — completeness and accuracy are critical. Missing even one entity, number, or relationship could break downstream analysis.

═══ PHASE 1: DOCUMENT STRUCTURE ANALYSIS ═══
Before extracting, identify the document type and layout:
- Document class (letter, report, invoice, form, slide, spreadsheet, article, legal, handwritten, mixed)
- Layout structure (single-column, multi-column, grid, free-form)
- Reading order (left-to-right, top-to-bottom, or complex multi-region)

═══ PHASE 2: FULL TEXT EXTRACTION ═══
Extract every piece of text exactly as written — character-perfect, preserving:

STRUCTURE & HIERARCHY:
  • Headings (H1/H2/H3) — reproduce with # / ## / ### markers
  • Paragraphs — maintain paragraph breaks as written
  • Bullet points & numbered lists — preserve nesting depth and numbering style
  • Indentation levels — use spaces to reflect visual hierarchy

TABLES (critical — extract with precision):
  • Reproduce ALL tables using | column | separators |
  • Include header rows with --- separator beneath
  • Capture merged cells — repeat content in each merged position
  • Preserve alignment (numeric columns right-aligned where visible)
  • Empty cells → leave blank between pipes, never skip

FORMS & KEY-VALUE PAIRS:
  • Field Label: Value — extract every pair, including empty/unchecked fields
  • Checkboxes: [x] checked, [ ] unchecked
  • Radio buttons: (●) selected, (○) unselected
  • Dropdown/selection: note the selected value

METADATA & PERIPHERAL TEXT:
  • Headers, footers, page numbers
  • Watermarks, stamps, confidentiality notices
  • Margin annotations, sticky notes
  • Barcodes/QR codes — note presence and any human-readable text beneath
  • Fine print, disclaimers, footnotes, endnotes (with reference numbers)

SPECIAL TEXT:
  • Handwritten text → extract with [handwritten] tag, mark uncertain chars as [?]
  • Struck-through text → ~~struck text~~
  • Underlined text → note as __underlined__
  • Bold text → **bold**
  • Text in colored highlights → {highlighted in yellow: text here}
  • Superscript/subscript → note with ^super^ or ~sub~
  • Mathematical formulas → LaTeX notation where possible
  • Code snippets → wrap in ```language blocks

MULTILINGUAL:
  • Extract ALL languages present — do not translate, keep original
  • Note the language in brackets if non-obvious: [Japanese] テキスト

═══ PHASE 3: VISUAL CONTENT EXTRACTION ═══
Every non-text visual element must be captured with full analytical detail:

CHARTS & GRAPHS:
  • Type: bar, line, pie, scatter, area, waterfall, Gantt, heatmap, etc.
  • Title, subtitle, axis labels (exact text)
  • ALL data values — read every bar height, line point, pie percentage
  • Legend items with their colors/patterns
  • Trend lines, annotations, callout boxes
  • Format: [Chart: {type} — Title: "{title}" | X-axis: {label} | Y-axis: {label} | Data: {series}]

DIAGRAMS & FLOWCHARTS:
  • Every node/box with its exact label text
  • Every connection/arrow with direction and any label on it
  • Decision points (diamond shapes) with yes/no paths
  • Swim lanes with lane headers
  • Format: [Diagram: {type} — Nodes: {list} | Flows: {A → B (label), B → C}]

ORGANIZATIONAL CHARTS:
  • Every person/role box with name, title, department
  • Reporting lines and hierarchy
  • Format: [OrgChart: {root} → {direct reports} → ...]

IMAGES & PHOTOGRAPHS:
  • Describe the subject, context, and any visible text/labels
  • Product photos: note model numbers, brand names, features visible
  • Screenshots: extract ALL UI text, buttons, menus, status indicators
  • Maps: note locations, labels, legends, scale
  • Format: [Image: {description} | Text visible: "{text}"]

LOGOS & BRANDING:
  • Company/organization name from logo
  • Any tagline or text within/near the logo
  • Format: [Logo: {company} — "{tagline}"]

SIGNATURES & STAMPS:
  • [Signature: {printed name if present, otherwise "illegible"}]
  • [Stamp: {text within stamp, date if visible}]

═══ PHASE 4: SPATIAL RELATIONSHIPS ═══
When content positioning matters for meaning:
  • Side-by-side comparisons → note left vs right
  • Callout arrows pointing from annotation to content
  • Grouped/boxed sections → note the grouping boundary and label
  • Color-coded regions → note what color means what

═══ OUTPUT RULES ═══
1. Output extracted text FIRST in document reading order, preserving all formatting
2. Visual elements inline where they appear in the document flow using [bracketed] notation
3. NEVER summarize, paraphrase, or interpret — extract VERBATIM
4. NEVER hallucinate content — if uncertain, use [unclear] or [partially illegible: best guess]
5. NEVER skip "unimportant" content — everything matters for the knowledge graph
6. If the image is blank or contains no discernible content, respond with exactly: [No text detected]
7. Maintain the EXACT same numbers, dates, currencies, percentages as shown — do not round or convert"""


def ocr_image_with_llm(image_bytes, mime_type='image/png', llm_inst=None):
    """Send an image to the LLM for text extraction via vision."""
    import base64
    _llm = llm_inst or llm
    completion = _llm.new_completion()
    completion.with_message(OCR_SYSTEM_PROMPT, role='system')
    mp = completion.new_multipart_message(role='user')
    mp.with_text(OCR_PROMPT)
    if isinstance(image_bytes, bytes):
        img_b64 = base64.b64encode(image_bytes).decode('utf-8')
    else:
        img_b64 = image_bytes
    mp.with_inline_image(img_b64, mime_type=mime_type)
    mp.add()
    resp = completion.execute()
    return resp.text.strip()


def extract_text_from_file(file_obj, filename, llm_inst=None, on_progress=None):
    import io
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    _progress = on_progress or (lambda msg, technique: None)

    # --- Image files: direct LLM OCR ---
    if ext in IMAGE_EXTENSIONS:
        _progress('Sending image to LLM Vision for text extraction...', 'LLM Vision OCR')
        image_bytes = file_obj.read()
        mime = IMAGE_MIME_MAP.get(ext, 'image/png')
        text = ocr_image_with_llm(image_bytes, mime_type=mime, llm_inst=llm_inst)
        _progress('Text extraction complete', 'LLM Vision OCR')
        return '[Image: ' + filename + ']\n\n' + text

    if ext in ('txt', 'md'):
        _progress('Reading plain text file...', 'Direct Read')
        return file_obj.read().decode('utf-8', errors='ignore')

    elif ext == 'csv':
        try:
            _progress('Parsing CSV structure...', 'CSV Parser')
            import csv
            raw = file_obj.read().decode('utf-8', errors='ignore')
            reader = csv.reader(raw.strip().splitlines())
            rows = list(reader)
            if not rows:
                return raw
            header = rows[0]
            parts = ['[Table: ' + filename + ', ' + str(len(rows)-1) + ' rows]', '']
            for ri, row in enumerate(rows[1:], 1):
                record = ['[Record ' + str(ri) + ']']
                for i, val in enumerate(row):
                    if val.strip():
                        col = header[i] if i < len(header) else 'Col' + str(i+1)
                        record.append('  ' + col + ': ' + val.strip())
                if len(record) > 1:
                    parts.append('\n'.join(record))
                    parts.append('')
            return '\n'.join(parts)
        except Exception:
            return file_obj.read().decode('utf-8', errors='ignore')

    elif ext == 'pdf':
        try:
            import fitz
            pdf_bytes = file_obj.read()
            pdf = fitz.open(stream=pdf_bytes, filetype='pdf')
            total_pages = len(pdf)
            _progress('Opened PDF — ' + str(total_pages) + ' pages. Using LLM Vision for full capture.', 'LLM Vision OCR')
            parts = ['[Document: ' + filename + ', ' + str(total_pages) + ' pages]', '']
            for i, page in enumerate(pdf):
                _progress('Page ' + str(i+1) + '/' + str(total_pages) + ' — rendering & sending to LLM Vision...', 'LLM Vision OCR')
                # Always render page as image so LLM sees text + embedded images/charts/diagrams
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes('png')
                try:
                    page_text = ocr_image_with_llm(img_bytes, mime_type='image/png', llm_inst=llm_inst)
                    if page_text and page_text != '[No text detected]':
                        parts.append('--- Page ' + str(i+1) + ' ---')
                        parts.append(page_text)
                        parts.append('')
                    _progress('Page ' + str(i+1) + '/' + str(total_pages) + ' — complete', 'LLM Vision OCR')
                except Exception as e:
                    parts.append('--- Page ' + str(i+1) + ' ---')
                    parts.append('[Vision failed: ' + str(e) + ']')
                    parts.append('')
                    _progress('Page ' + str(i+1) + ' — failed: ' + str(e), 'LLM Vision OCR')
            _progress('All ' + str(total_pages) + ' pages processed via LLM Vision', 'Done')
            return '\n'.join(parts)
        except ImportError:
            return '[Error: PyMuPDF not installed for PDF support]'

    elif ext == 'docx':
        try:
            _progress('Parsing Word document...', 'python-docx')
            from docx import Document
            doc = Document(io.BytesIO(file_obj.read()))
            parts = ['[Document: ' + filename + ']', '']
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style = para.style.name if para.style else ''
                if 'Heading 1' in style:
                    parts.append('# ' + text)
                elif 'Heading 2' in style:
                    parts.append('## ' + text)
                elif 'Heading 3' in style:
                    parts.append('### ' + text)
                elif 'List' in style:
                    parts.append('- ' + text)
                else:
                    parts.append(text)
            for table in doc.tables:
                parts.append('')
                parts.append('[Table]')
                for ri, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(' | '.join(cells))
                    if ri == 0:
                        parts.append('-' * 40)
                parts.append('')
            return '\n'.join(parts)
        except ImportError:
            return '[Error: python-docx not installed for DOCX support]'

    elif ext == 'pptx':
        pptx_bytes = file_obj.read()
        # Try rendering slides as images via PyMuPDF for full LLM Vision extraction
        try:
            import fitz
            doc = fitz.open(stream=pptx_bytes, filetype='pptx')
            total_slides = len(doc)
            _progress('Opened PPTX — ' + str(total_slides) + ' slides. Using LLM Vision for full capture.', 'LLM Vision')
            parts = ['[Presentation: ' + filename + ', ' + str(total_slides) + ' slides]', '']
            # Also extract speaker notes via python-pptx (vision can't see these)
            slide_notes = {}
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(pptx_bytes))
                for ni, slide in enumerate(prs.slides):
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            slide_notes[ni] = notes
            except Exception:
                pass
            for i, page in enumerate(doc):
                _progress('Slide ' + str(i+1) + '/' + str(total_slides) + ' — rendering & sending to LLM Vision...', 'LLM Vision')
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes('png')
                try:
                    slide_text = ocr_image_with_llm(img_bytes, mime_type='image/png', llm_inst=llm_inst)
                    if slide_text and slide_text != '[No text detected]':
                        parts.append('--- Slide ' + str(i+1) + ' ---')
                        parts.append(slide_text)
                        if i in slide_notes:
                            parts.append('[Speaker Notes] ' + slide_notes[i])
                        parts.append('')
                    _progress('Slide ' + str(i+1) + '/' + str(total_slides) + ' — complete', 'LLM Vision')
                except Exception as e:
                    parts.append('--- Slide ' + str(i+1) + ' ---')
                    parts.append('[Vision failed: ' + str(e) + ']')
                    parts.append('')
                    _progress('Slide ' + str(i+1) + ' — failed: ' + str(e), 'LLM Vision')
            _progress('All ' + str(total_slides) + ' slides processed via LLM Vision', 'Done')
            return '\n'.join(parts)
        except (ImportError, Exception):
            # Fallback to text-only extraction if PyMuPDF can't render PPTX
            try:
                _progress('PyMuPDF unavailable for PPTX rendering — falling back to text extraction...', 'python-pptx')
                from pptx import Presentation
                prs = Presentation(io.BytesIO(pptx_bytes))
                parts = ['[Presentation: ' + filename + ', ' + str(len(prs.slides)) + ' slides]', '']
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    slide_texts.append(t)
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                slide_texts.append(' | '.join(cells))
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            slide_texts.append('[Speaker Notes] ' + notes)
                    if slide_texts:
                        parts.append('--- Slide ' + str(i+1) + ' ---')
                        parts.extend(slide_texts)
                        parts.append('')
                return '\n'.join(parts)
            except ImportError:
                return '[Error: python-pptx not installed for PPTX support]'

    elif ext == 'xlsx':
        try:
            _progress('Parsing Excel spreadsheet...', 'openpyxl')
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_obj.read()), read_only=True)
            parts = ['[Spreadsheet: ' + filename + ', ' + str(len(wb.sheetnames)) + ' sheets]', '']
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                parts.append('--- Sheet: ' + ws.title + ' ---')
                header = [str(c) if c is not None else '' for c in rows[0]]
                parts.append('')
                for ri, row in enumerate(rows[1:], 1):
                    record = ['[Record ' + str(ri) + ']']
                    for i, val in enumerate(row):
                        if val is not None and str(val).strip():
                            col = header[i] if i < len(header) and header[i] else 'Col' + str(i+1)
                            record.append('  ' + col + ': ' + str(val).strip())
                    if len(record) > 1:
                        parts.append('\n'.join(record))
                        parts.append('')
            return '\n'.join(parts)
        except ImportError:
            return '[Error: openpyxl not installed for XLSX support]'

    elif ext == 'eml':
        try:
            _progress('Parsing email (.eml)...', 'Email Parser')
            import email
            from email import policy
            msg = email.message_from_bytes(file_obj.read(), policy=policy.default)
            parts = ['[Email: ' + filename + ']']
            parts.append('From: ' + str(msg.get('From', '')))
            parts.append('To: ' + str(msg.get('To', '')))
            if msg.get('Cc'):
                parts.append('Cc: ' + str(msg.get('Cc', '')))
            parts.append('Date: ' + str(msg.get('Date', '')))
            parts.append('Subject: ' + str(msg.get('Subject', '')))
            parts.append('')
            body = msg.get_body(preferencelist=('plain', 'html'))
            if body:
                content = body.get_content()
                if body.get_content_type() == 'text/html':
                    import re as _re
                    content = _re.sub(r'<[^>]+>', ' ', content)
                    content = _re.sub(r'\s+', ' ', content).strip()
                parts.append(content)
            return '\n'.join(parts)
        except Exception as e:
            return '[Error parsing EML: ' + str(e) + ']'

    elif ext == 'msg':
        try:
            _progress('Parsing Outlook email (.msg)...', 'extract-msg')
            import extract_msg
            msg = extract_msg.Message(io.BytesIO(file_obj.read()))
            parts = ['[Email: ' + filename + ']']
            parts.append('From: ' + str(msg.sender or ''))
            parts.append('To: ' + str(msg.to or ''))
            parts.append('Date: ' + str(msg.date or ''))
            parts.append('Subject: ' + str(msg.subject or ''))
            parts.append('')
            parts.append(str(msg.body or ''))
            return '\n'.join(parts)
        except ImportError:
            return '[Error: extract-msg not installed for MSG support]'

    else:
        return file_obj.read().decode('utf-8', errors='ignore')


# --- Flask Routes ---
@app.route('/models')
def api_models():
    models = list_available_llms()
    return json_response({'models': models})


@app.route('/domains')
def api_domains():
    domains = [{'key': 'auto', 'name': 'Auto-detect'}]
    for key, tmpl in DOMAIN_TEMPLATES.items():
        if key == 'general':
            domains.append({'key': 'general', 'name': 'General'})
        else:
            domains.append({'key': key, 'name': tmpl['name']})
    return json_response({'domains': domains})


@app.route('/upload', methods=['POST'])
def api_upload():
    if 'file' not in request.files:
        return json_response({'error': 'No file provided'}, 400)
    f = request.files['file']
    if not f.filename:
        return json_response({'error': 'No file selected'}, 400)
    model_id = request.form.get('model', None)
    llm_inst = get_llm(model_id) if model_id else llm
    fname = f.filename

    # Read file into memory so we can stream response
    import io
    file_bytes = f.read()
    file_like = io.BytesIO(file_bytes)

    # Collect progress events and final text via a queue
    import queue
    q = queue.Queue()

    def on_progress(msg, technique):
        q.put(('progress', msg, technique))

    def run_extract():
        try:
            text = extract_text_from_file(file_like, fname, llm_inst=llm_inst, on_progress=on_progress)
            q.put(('done', text, fname))
        except Exception as e:
            q.put(('error', str(e), ''))

    t = threading.Thread(target=run_extract)
    t.start()

    def generate():
        while True:
            try:
                item = q.get(timeout=120)
            except Exception:
                yield 'data: ' + json.dumps({'error': 'Timeout'}) + '\n\n'
                break
            if item[0] == 'progress':
                yield 'data: ' + json.dumps({'progress': item[1], 'technique': item[2]}) + '\n\n'
            elif item[0] == 'done':
                yield 'data: ' + json.dumps({'done': True, 'text': item[1], 'filename': item[2]}) + '\n\n'
                break
            elif item[0] == 'error':
                yield 'data: ' + json.dumps({'error': item[1]}) + '\n\n'
                break

    return Response(generate(), mimetype='text/event-stream',
                    headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'})


@app.route('/build', methods=['POST'])
def api_build():
    data = request.get_json(force=True)
    text = data.get('text', '').strip()
    if not text:
        return json_response({'error': 'No text provided'}, 400)

    session_id = str(uuid.uuid4())
    sessions[session_id] = {
        'status': 'starting',
        'current_chunk': 0,
        'total_chunks': 0,
        'graph': None,
        'graph_data': {'nodes': [], 'edges': [], 'typeColors': {}},
        'source_text': text,
        'messages': [],
        'errors': []
    }

    model_id = data.get('model', None)
    domain = data.get('domain', 'auto')
    thread = threading.Thread(target=build_graph_async, args=(session_id, text, model_id, domain))
    thread.daemon = True
    thread.start()

    return json_response({'session_id': session_id})


@app.route('/status/<session_id>')
def api_status(session_id):
    session = sessions.get(session_id)
    if not session:
        return json_response({'error': 'Session not found'}, 404)
    return json_response({
        'status': session['status'],
        'current_chunk': session['current_chunk'],
        'total_chunks': session['total_chunks'],
        'graph_data': session['graph_data'],
        'errors': session.get('errors', []),
        'schema': session.get('schema', None),
        'summary': session.get('summary', None),
        'report': session.get('report', None)
    })


@app.route('/stop/<session_id>', methods=['POST'])
def api_stop(session_id):
    session = sessions.get(session_id)
    if not session:
        return json_response({'error': 'Session not found'}, 404)
    session['stopped'] = True
    return json_response({'ok': True})


@app.route('/report/<session_id>', methods=['POST'])
def api_report(session_id):
    session = sessions.get(session_id)
    if not session or not session.get('graph'):
        return json_response({'error': 'No graph built yet'}, 400)

    # Return cached report if already generated
    if session.get('report'):
        return json_response({'report': session['report']})

    # Otherwise generate fresh
    _llm = session.get('llm_inst', llm)
    report = generate_report(session, _llm)
    session['report'] = report
    return json_response({'report': report})


MAX_HISTORY_TURNS = 10  # Keep last N Q&A pairs for context
MAX_CONTEXT_CHARS = 6000  # Max chars for source text in context

def build_qa_system_prompt(session):
    """Build the system/context prompt once per session, with truncated source."""
    ctx = get_full_context(session['graph'])
    source = session['source_text']
    if len(source) > MAX_CONTEXT_CHARS:
        source = source[:MAX_CONTEXT_CHARS] + '\n[... truncated ...]'
    return """You are a research assistant for the Orion platform. Answer questions using ONLY the provided data.
If the answer is not in the data, say so. Be concise and professional.
When referencing entities or relationships, be specific about names and types.
If a follow-up question references something from the conversation, use context from previous exchanges.

Entity & Relationship Data:
""" + ctx + """

Source Text:
""" + source


@app.route('/ask/<session_id>', methods=['POST'])
def api_ask(session_id):
    session = sessions.get(session_id)
    if not session or not session.get('graph'):
        return json_response({'error': 'No graph built yet'}, 400)

    data = request.get_json(force=True)
    question = data.get('question', '').strip()
    if not question:
        return json_response({'error': 'No question provided'}, 400)

    _llm = session.get('llm_inst', llm)
    completion = _llm.new_completion()

    # System context with entity data
    system_prompt = build_qa_system_prompt(session)
    completion.with_message(system_prompt, role='system')

    # Sliding window: include last N turns of conversation history
    history = session.get('messages', [])
    window = history[-(MAX_HISTORY_TURNS * 2):]  # Each turn = 2 messages (user + assistant)
    for msg in window:
        completion.with_message(msg['content'], role=msg['role'])

    # Current question
    completion.with_message(question, role='user')
    resp = completion.execute()

    # Append to history
    session['messages'].append({'role': 'user', 'content': question})
    session['messages'].append({'role': 'assistant', 'content': resp.text})

    return json_response({'answer': resp.text})


@app.route('/ask_stream/<session_id>', methods=['POST'])
def api_ask_stream(session_id):
    """Streaming Q&A endpoint — returns SSE events word-by-word."""
    import time
    session = sessions.get(session_id)
    if not session or not session.get('graph'):
        return json_response({'error': 'No graph built yet'}, 400)

    data = request.get_json(force=True)
    question = data.get('question', '').strip()
    if not question:
        return json_response({'error': 'No question provided'}, 400)

    _llm = session.get('llm_inst', llm)
    completion = _llm.new_completion()

    system_prompt = build_qa_system_prompt(session)
    completion.with_message(system_prompt, role='system')

    history = session.get('messages', [])
    window = history[-(MAX_HISTORY_TURNS * 2):]
    for msg in window:
        completion.with_message(msg['content'], role=msg['role'])

    completion.with_message(question, role='user')
    resp = completion.execute()
    answer = resp.text

    # Save to history
    session['messages'].append({'role': 'user', 'content': question})
    session['messages'].append({'role': 'assistant', 'content': answer})

    def generate():
        words = answer.split(' ')
        for i, word in enumerate(words):
            chunk = word + (' ' if i < len(words) - 1 else '')
            yield 'data: ' + json.dumps({'chunk': chunk}) + '\n\n'
            time.sleep(0.03)  # ~30ms per word for natural feel
        yield 'data: ' + json.dumps({'done': True}) + '\n\n'

    return Response(generate(), mimetype='text/event-stream',
                    headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'})
